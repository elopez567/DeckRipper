import pyscreenshot as ImageGrab
import time
from selenium import webdriver
import tkinter as tk
from tkinter import *
from PIL import Image, ImageTk, ImageGrab
from selenium.webdriver.common.by import By
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx import Presentation
from pptx.util import Inches
import pptx

driver = webdriver.Chrome('C:/Users/elopez1/PythonProjects/Roadshow SS/chromedriver')
driver.get('https://dealroadshow.com')
time.sleep(1)
email_input = driver.find_element(By.XPATH,
                                          '/html/body/div[1]/div[2]/div/div[1]/div[1]/div/div[1]/form/div[1]/label[2]/input')
email_input.send_keys('Emmanuel.Lopez@pnmac.com')


class Mainapp:
    def __init__(self):
        dim = None
        app = Tk()
        app.geometry("250x500")
        slide_count_label = Label(app, text='                        # of Slides', font=('Calibri', 11), pady=20)
        slide_count_label.grid(row=0, column=0, sticky=W)

        # Entry
        self.slide_count_text = IntVar()
        Entry(app, textvariable=self.slide_count_text).grid(row=1, column=1, pady=20)

        Button(app, text='Rip Deck', width=15,
               command=self.Screenshot).grid(row=4, column=1, pady=20)

        Button(app, text='Select Area', width=15,
               command=ScreenSize).grid(row=2, column=1, pady=20)

        Button(app, text='Record Dimensions', width=15,
               command=dimensions).grid(row=3, column=1, pady=20)

        Button(app, text='Create PP', width=15,
               command=self.CreatePP).grid(row=5, column=1, pady=20)

        app.mainloop()

    def Screenshot(self):
        slide_count = self.slide_count_text.get()
        next_button = driver.find_element(By.XPATH,
                                          '/html/body/div[1]/div/div/div/div/div[2]/div/div[1]/div[3]/div/div[2]/span/span/button')
        for i in range(slide_count):
            im = ImageGrab.grab(bbox=(x1, y1, x2, y2))
            im.save(f'C:/Users/elopez1/PythonProjects/Roadshow SS/Screenshots/Slide {i}.png')
            next_button.click()
            time.sleep(1)

    def CreatePP(self):
        slide_count = self.slide_count_text.get()
        prs = Presentation()
        layout = prs.slide_layouts[6]

        for i in range(slide_count):
            slide = prs.slides.add_slide(layout)
            img_path = f'C:/Users/elopez1/PythonProjects/Roadshow SS/Screenshots/Slide {i}.png'

            im = Image.open(img_path)
            height = im.height
            pic = slide.shapes.add_picture(img_path, pptx.util.Inches(0), pptx.util.Inches(1),
                                           width=pptx.util.Inches(10))
            pic.top = int(height)

        prs.save('C:/Users/elopez1/PythonProjects/Roadshow SS/test.pptx')




class ScreenSize:
    def __init__(self):
        WIDTH, HEIGHT = 900, 900
        topx, topy, botx, boty = 0, 0, 0, 0
        global canvas, rect_id, window
        rect_id = None
        path = 'C:/Users/elopez1/PythonProjects/Roadshow SS/Screenshots/FullScreen.png'

        im=ImageGrab.grab()
        im.save("C:/Users/elopez1/PythonProjects/Roadshow SS/Screenshots/FullScreen.png")




        window = Toplevel()
        window.title("Select Area")
        window.geometry('%sx%s' % (WIDTH, HEIGHT))
        window.configure(background='grey')

        img = ImageTk.PhotoImage(Image.open(path))


        canvas = tk.Canvas(window, width=img.width(), height=img.height(),
                           borderwidth=0, highlightthickness=0)

        PhotoImage(master=canvas, width=WIDTH, height=HEIGHT)

        canvas.pack(expand=True)
        canvas.img = img  # Keep reference in case this code is put into a function.
        canvas.create_image(0, 0, image=img, anchor=tk.NW)

        # Create selection rectangle (invisible since corner points are equal).
        rect_id = canvas.create_rectangle(topx, topy, topx, topy,
                                          dash=(2,2), fill='', outline='white')
        def get_mouse_posn(event):
            global topy, topx

            topx, topy = event.x, event.y

        def update_sel_rect(event):
            global rect_id, canvas
            global topy, topx, botx, boty

            botx, boty = event.x, event.y
            canvas.coords(rect_id, topx, topy, botx, boty)  # Update selection rect.


        canvas.bind('<Button-1>', get_mouse_posn)
        canvas.bind('<B1-Motion>', update_sel_rect)


def dimensions():
    dlist = canvas.coords(rect_id)
    global x1, y1, x2, y2
    x1, y1, x2, y2 = [i for i in dlist]
    print(dlist)
    window.destroy()


Mainapp()
